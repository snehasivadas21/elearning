import requests
from io import BytesIO
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from ai.vectorstore import get_vectorstore
from ai.models import CourseEmbedding
from django.utils import timezone
from courses.models import LessonResource


def extract_text(resource: LessonResource):
    url = resource.file.url
    response = requests.get(url, timeout=30)
    response.raise_for_status()

    content_type = response.headers.get("Content-Type", "").lower()
    content = BytesIO(response.content)

    if "pdf" in content_type:
        reader = PdfReader(content)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if "word" in content_type or "docx" in content_type:
        doc = Document(content)
        return "\n".join(p.text for p in doc.paragraphs)

    if "presentation" in content_type or "ppt" in content_type:
        prs = Presentation(content)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    return ""

def index_lesson_resource(resource):
    text = extract_text(resource)

    if not text.strip():
        print("⚠️ No text extracted for:", resource.id)
        return

    vectorstore = get_vectorstore()

    metadata = {
        "course_id": str(resource.lesson.module.course.id),
        "lesson_id": str(resource.lesson.id),
        "resource_id": str(resource.id),
        "source_type": "lesson_resources",
    }

    vectorstore.add_texts(
        texts=[text],
        metadatas=[metadata]
    )

    CourseEmbedding.objects.create(
        course=resource.lesson.module.course,
        lesson=resource.lesson,
        content=text[:10000],
        source_type="lesson_resources",
        is_indexed=True,
        indexed_at=timezone.now(),
    )

    print("✅ Indexed resource:", resource.id)
